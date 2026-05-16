import os
import gc
import tempfile
import cv2
import numpy as np
from utils import get_poppler_path, get_model_path

def process_pdf_to_word(input_file, output_path, status_callback, stop_event):
    from pdf2docx import Converter
    status_callback("📝 啟動引擎，正在轉換 PDF 為 Word...", 0.5)
    cv = Converter(input_file)
    cv.convert(output_path, start=0, end=None)
    cv.close()

def process_pdf_to_excel(input_file, output_path, status_callback, stop_event):
    import pdfplumber
    import pandas as pd
    status_callback("📊 正在提取表格資料...", 0.2)
    all_tables = []
    with pdfplumber.open(input_file) as pdf:
        total = len(pdf.pages)
        for i, page in enumerate(pdf.pages):
            if stop_event.is_set(): return
            status_callback(f"📊 掃描第 {i+1} 頁表格...", (i+1)/total)
            tables = page.extract_tables()
            for table in tables:
                if table and len(table) > 1:
                    df = pd.DataFrame(table[1:], columns=table[0])
                    all_tables.append(df)
                    
    if not all_tables: raise Exception("此 PDF 中找不到任何表格資料！")
    
    status_callback("💾 正在寫入 Excel...", 0.9)
    with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
        for i, df in enumerate(all_tables):
            df.to_excel(writer, sheet_name=f"Table_{i+1}", index=False)

def process_pdf_to_images(input_file, output_dir, status_callback, stop_event, dpi=300):
    from pdf2image import convert_from_path, pdfinfo_from_path
    poppler = get_poppler_path()
    info = pdfinfo_from_path(input_file, poppler_path=poppler)
    total = info["Pages"]
    base_name = os.path.splitext(os.path.basename(input_file))[0]
    for i in range(1, total + 1):
        if stop_event.is_set(): break
        status_callback(f"🖼️ 正在處理並儲存圖片 {i} / {total}...", i/total)
        page_img = convert_from_path(input_file, dpi=dpi, first_page=i, last_page=i, poppler_path=poppler)[0]
        page_img.save(os.path.join(output_dir, f"{base_name}_{i}.jpg"), 'JPEG')
        del page_img; gc.collect()

def process_images_to_pdf(input_files, output_path, status_callback, stop_event):
    import fitz  
    doc = fitz.open()
    total = len(input_files)
    for i, img_path in enumerate(input_files):
        if stop_event.is_set(): return
        status_callback(f"🖼️ 正在將圖片轉為 PDF... ({i+1}/{total})", (i+1)/total)
        img_doc = fitz.open(img_path)
        pdf_bytes = img_doc.convert_to_pdf()
        img_pdf = fitz.open("pdf", pdf_bytes)
        doc.insert_pdf(img_pdf)
        img_doc.close(); img_pdf.close()
        del img_doc; gc.collect()
    doc.save(output_path)
    doc.close()

def process_pdf_to_ppt(input_file, output_path, status_callback, stop_event, dpi=300, ppt_mode="圖文排版 (智慧 OCR)"):
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    import fitz
    from PIL import Image, ImageDraw
    from collections import Counter
    import opencc
    from bs4 import BeautifulSoup
    
    converter = opencc.OpenCC('s2twp')
    is_pdf = input_file.lower().endswith('.pdf')
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6] 
    
    def expand_bbox(x0, y0, x1, y1, scale=1.1, max_w=9999, max_h=9999):
        cx, cy = (x0 + x1) / 2, (y0 + y1) / 2
        w, h = (x1 - x0) * scale, (y1 - y0) * scale
        return [max(0, cx - w / 2), max(0, cy - h / 2), min(max_w, cx + w / 2), min(max_h, cy + h / 2)]

    def get_dynamic_bg_color(img_obj, px_bbox):
        w, h = img_obj.size
        x0, y0, x1, y1 = [int(v) for v in px_bbox]
        samples = [(x0-5, y0-5), (x0+(x1-x0)//2, y0-5), (x1+5, y0-5), (x0-5, y1+5), (x0+(x1-x0)//2, y1+5), (x1+5, y1+5), (x0-5, y0+(y1-y0)//2), (x1+5, y0+(y1-y0)//2)]
        valid_samples = []
        for sx, sy in samples:
            try:
                p = img_obj.getpixel((max(0, min(w-1, sx)), max(0, min(h-1, sy))))
                valid_samples.append(p[:3] if isinstance(p, tuple) else (p,p,p))
            except: pass
        return Counter(valid_samples).most_common(1)[0][0] if valid_samples else (255, 255, 255)

    def draw_ppt_table_from_html(slide, html_str, x_pt, y_pt, w_pt, h_pt):
        soup = BeautifulSoup(html_str, 'html.parser')
        rows = soup.find_all('tr')
        if not rows: return
        num_rows = len(rows)
        num_cols = max(len(row.find_all(['td', 'th'])) for row in rows)
        table = slide.shapes.add_table(num_rows, num_cols, Pt(x_pt), Pt(y_pt), Pt(w_pt), Pt(h_pt)).table
        for r_idx, row in enumerate(rows):
            for c_idx, col in enumerate(row.find_all(['td', 'th'])):
                if r_idx < len(table.rows) and c_idx < len(table.columns):
                    cell = table.cell(r_idx, c_idx)
                    cell.text = converter.convert(col.get_text(strip=True))
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.name = "微軟正黑體"
                            run.font.size = Pt(10)
                            run.font.color.rgb = RGBColor(0, 0, 0)

    ocr_engine, layout_engine, table_engine = None, None, None
    if ppt_mode == "圖文排版 (智慧 OCR)":
        from rapidocr_onnxruntime import RapidOCR
        from rapid_layout import RapidLayout
        from rapid_table import RapidTable
        model_dir = get_model_path()
        det_path = os.path.join(model_dir, "ch_PP-OCRv4_det_infer.onnx")
        if os.path.exists(det_path):
            ocr_engine = RapidOCR(det_model_path=det_path, cls_model_path=os.path.join(model_dir, "ch_ppocr_mobile_v2.0_cls_infer.onnx"), rec_model_path=os.path.join(model_dir, "ch_PP-OCRv4_rec_infer.onnx"), use_angle_cls=True, det_box_thresh=0.2, text_score=0.2)
            layout_engine = RapidLayout(model_path=os.path.join(model_dir, "layout_cdla_dataset.onnx"))
            table_engine = RapidTable(model_path=os.path.join(model_dir, "ch_ppstructure_mobile_v2.0_SLANet.onnx"))
        else:
            ocr_engine = RapidOCR(use_angle_cls=True, det_box_thresh=0.2, text_score=0.2)
            layout_engine = RapidLayout()
            table_engine = RapidTable()

    with tempfile.TemporaryDirectory() as temp_dir:
        def process_slide(hr_img_path, normal_img_path, slide_obj, scale_down):
            img_obj = Image.open(normal_img_path).convert("RGB")
            draw = ImageDraw.Draw(img_obj)
            text_boxes_data, table_html_data = [], []
            hr_img = cv2.imdecode(np.fromfile(hr_img_path, dtype=np.uint8), cv2.IMREAD_COLOR)
            if hr_img is None: return 
            
            def enhance_for_ocr(image):
                try:
                    img_yuv = cv2.cvtColor(image, cv2.COLOR_BGR2YUV)
                    img_yuv[:,:,0] = cv2.createCLAHE(clipLimit=1.5, tileGridSize=(8, 8)).apply(img_yuv[:,:,0])
                    return cv2.cvtColor(img_yuv, cv2.COLOR_YUV2BGR)
                except: return image 

            raw_layout = layout_engine(hr_img)
            regions = []
            if isinstance(raw_layout, tuple): raw_layout = raw_layout[0]
            if isinstance(raw_layout, list):
                for item in raw_layout:
                    if isinstance(item, dict): regions.append((item.get('bbox'), item.get('label')))
            elif hasattr(raw_layout, 'boxes'):
                for box, label in zip(raw_layout.boxes, raw_layout.class_names): regions.append((box, label))

            if regions:
                for box, label in regions:
                    if not box or not label: continue
                    if len(box) == 4 and not isinstance(box[0], (list, tuple, np.ndarray)): bx0, by0, bx1, by1 = box
                    else: bx0, by0, bx1, by1 = min([p[0] for p in box]), min([p[1] for p in box]), max([p[0] for p in box]), max([p[1] for p in box])
                    x0, y0, x1, y1 = [v * scale_down for v in [bx0, by0, bx1, by1]]
                    
                    if label == 'table':
                        table_crop = hr_img[int(by0):int(by1), int(bx0):int(bx1)]
                        if table_crop is not None and table_crop.size > 0:
                            table_res = table_engine(enhance_for_ocr(table_crop))
                            table_res = table_res[0] if isinstance(table_res, tuple) else table_res
                            html_str = table_res if isinstance(table_res, str) else table_res.get('html') if isinstance(table_res, dict) else getattr(table_res, 'html', None)
                            if html_str:
                                table_html_data.append({"html": html_str, "x": x0, "y": y0, "w": x1-x0, "h": y1-y0})
                                exp_px_bbox = expand_bbox(x0, y0, x1, y1, scale=1.02, max_w=img_obj.width, max_h=img_obj.height)
                                draw.rectangle(exp_px_bbox, fill=get_dynamic_bg_color(img_obj, exp_px_bbox))
                    elif label in ['text', 'title', 'figure']:
                        text_crop = hr_img[int(by0):int(by1), int(bx0):int(bx1)]
                        if text_crop is not None and text_crop.size > 0:
                            ocr_res = ocr_engine(enhance_for_ocr(text_crop))
                            ocr_res = ocr_res[0] if isinstance(ocr_res, tuple) else ocr_res
                            if ocr_res:
                                for line in ocr_res:
                                    lx0, ly0, lx1, ly1 = min([p[0] for p in line[0]])*scale_down+x0, min([p[1] for p in line[0]])*scale_down+y0, max([p[0] for p in line[0]])*scale_down+x0, max([p[1] for p in line[0]])*scale_down+y0
                                    text_boxes_data.append({"text": converter.convert(line[1]), "x": lx0, "y": ly0, "w": lx1-lx0, "h": ly1-ly0})
                                    exp_px_bbox = expand_bbox(lx0, ly0, lx1, ly1, scale=1.1, max_w=img_obj.width, max_h=img_obj.height)
                                    draw.rectangle(exp_px_bbox, fill=get_dynamic_bg_color(img_obj, exp_px_bbox))
            else:
                ocr_res = ocr_engine(enhance_for_ocr(hr_img))
                ocr_res = ocr_res[0] if isinstance(ocr_res, tuple) else ocr_res
                if ocr_res:
                    for line in ocr_res:
                        lx0, ly0, lx1, ly1 = min([p[0] for p in line[0]])*scale_down, min([p[1] for p in line[0]])*scale_down, max([p[0] for p in line[0]])*scale_down, max([p[1] for p in line[0]])*scale_down
                        text_boxes_data.append({"text": converter.convert(line[1]), "x": lx0, "y": ly0, "w": lx1-lx0, "h": ly1-ly0})
                        exp_px_bbox = expand_bbox(lx0, ly0, lx1, ly1, scale=1.1, max_w=img_obj.width, max_h=img_obj.height)
                        draw.rectangle(exp_px_bbox, fill=get_dynamic_bg_color(img_obj, exp_px_bbox))

            img_obj.save(normal_img_path, "JPEG", quality=90)
            slide_obj.shapes.add_picture(normal_img_path, 0, 0, prs.slide_width, prs.slide_height)
            for item in text_boxes_data:
                tf = slide_obj.shapes.add_textbox(Pt(item["x"]), Pt(item["y"]), Pt(item["w"]), Pt(item["h"])).text_frame
                tf.clear(); tf.word_wrap = False
                run = tf.paragraphs[0].add_run()
                run.text, run.font.size, run.font.name, run.font.color.rgb = item["text"], Pt(max(8, item["h"] * 0.75)), "微軟正黑體", RGBColor(0, 0, 0)
            for t_item in table_html_data: draw_ppt_table_from_html(slide_obj, t_item["html"], t_item["x"], t_item["y"], t_item["w"], t_item["h"])

        if is_pdf:
            doc = fitz.open(input_file)
            total = len(doc)
            for i in range(total):
                if stop_event.is_set(): break
                if i == 0: prs.slide_width, prs.slide_height = Pt(doc[i].rect.width), Pt(doc[i].rect.height)
                slide = prs.slides.add_slide(blank_slide_layout)
                normal_path = os.path.join(temp_dir, f"bg_{i}.jpg")
                if ppt_mode == "純圖片簡報 (較快)":
                    status_callback(f"📊 正在轉換純圖 PPT (第 {i+1} / {total} 頁)...", (i+1)/total)
                    doc[i].get_pixmap(dpi=dpi, colorspace=fitz.csRGB).save(normal_path)
                    slide.shapes.add_picture(normal_path, 0, 0, prs.slide_width, prs.slide_height)
                else:
                    status_callback(f"👁️ 版面分析與雙軌重建 (第 {i+1} / {total} 頁)...", (i+1)/total)
                    hr_path = os.path.join(temp_dir, f"hr_{i}.jpg")
                    doc[i].get_pixmap(matrix=fitz.Matrix(4.0, 4.0), colorspace=fitz.csRGB).save(hr_path)
                    doc[i].get_pixmap(dpi=72, colorspace=fitz.csRGB).save(normal_path)
                    process_slide(hr_path, normal_path, slide, scale_down=1.0/4.0)
            doc.close()
        else:
            img = Image.open(input_file).convert('RGB')
            normal_path, hr_path = os.path.join(temp_dir, "normal.jpg"), os.path.join(temp_dir, "hr.jpg")
            img.save(normal_path, "JPEG", quality=95)
            img.resize((img.width * 4, img.height * 4), Image.LANCZOS).save(hr_path, "JPEG", quality=100)
            prs.slide_width, prs.slide_height = Pt(img.width * 72 / dpi), Pt(img.height * 72 / dpi)
            slide = prs.slides.add_slide(blank_slide_layout)
            if ppt_mode == "純圖片簡報 (較快)":
                status_callback("🖼️ 正在處理圖片檔案...", 0.5)
                slide.shapes.add_picture(normal_path, 0, 0, prs.slide_width, prs.slide_height)
            else:
                status_callback("👁️ 進行版面分析與雙軌重建...", 0.5)
                process_slide(hr_path, normal_path, slide, scale_down=1.0/4.0)

        if not stop_event.is_set():
            status_callback("💾 正在儲存 PPT 檔案...", 0.95)
            prs.save(output_path)
