import os
import gc
import tempfile
import cv2
import numpy as np
from utils import get_poppler_path, apply_watermark_removal

def process_extract_text(input_file, output_path, status_callback, stop_event):
    import fitz
    doc = fitz.open(input_file)
    total = len(doc)
    with open(output_path, "w", encoding="utf-8") as f:
        for i, page in enumerate(doc):
            if stop_event.is_set(): return
            status_callback(f"📄 正在提取文字 ({i+1}/{total})...", (i+1)/total)
            f.write(f"--- 第 {i+1} 頁 ---\n{page.get_text('text')}\n\n")
    doc.close()

def process_extract_original_images(input_file, output_dir, status_callback, stop_event):
    import fitz
    doc = fitz.open(input_file)
    total = len(doc)
    count = 0
    for i in range(total):
        if stop_event.is_set(): return
        status_callback(f"🖼️ 正在提取內嵌圖片 (掃描第 {i+1} 頁)...", (i+1)/total)
        for img_idx, img in enumerate(doc[i].get_images(True)):
            base_img = doc.extract_image(img[0])
            with open(os.path.join(output_dir, f"page{i+1}_img{img_idx+1}.{base_img['ext']}"), "wb") as f:
                f.write(base_img["image"])
            count += 1
    doc.close()
    return f"提取完成！共成功提取了 {count} 張原始圖片。"

def process_add_watermark(input_file, output_path, text, status_callback, stop_event):
    import fitz
    doc = fitz.open(input_file)
    total = len(doc)
    for i, page in enumerate(doc):
        if stop_event.is_set(): return
        status_callback(f"🖋️ 正在加入文字浮水印... ({i+1}/{total})", (i+1)/total)
        page.insert_text(fitz.Point(50, page.rect.height / 2), text, fontsize=48, color=(1, 0, 0))
    doc.save(output_path)
    doc.close()

def process_add_image_watermark(input_file, output_path, img_path, position, target_page="全部頁面", status_callback=None, stop_event=None):
    import fitz
    doc = fitz.open(input_file)
    total = len(doc)
    img_doc = fitz.open(img_path)
    img_rect = img_doc[0].rect
    img_doc.close()
    
    target_indices = [0] if target_page == "僅第一頁" else [total - 1] if target_page == "僅最後一頁" else range(total)

    for i in target_indices:
        if stop_event and stop_event.is_set(): return
        if status_callback: status_callback(f"🖼️ 正在壓印圖片... (第 {i+1} 頁)", (i+1)/total)
        
        page = doc[i]
        page_rect = page.rect
        w = page_rect.width * 0.25
        h = w * (img_rect.height / img_rect.width)
        
        if position == "右下角": target_rect = fitz.Rect(page_rect.width - w - 20, page_rect.height - h - 20, page_rect.width - 20, page_rect.height - 20)
        elif position == "左下角": target_rect = fitz.Rect(20, page_rect.height - h - 20, 20 + w, page_rect.height - 20)
        elif position == "右上角": target_rect = fitz.Rect(page_rect.width - w - 20, 20, page_rect.width - 20, 20 + h)
        elif position == "正中央": target_rect = fitz.Rect(page_rect.width/2 - w/2, page_rect.height/2 - h/2, page_rect.width/2 + w/2, page_rect.height/2 + h/2)
        else: target_rect = fitz.Rect(20, 20, 20 + w, 20 + h)
        
        page.insert_image(target_rect, filename=img_path)
    
    doc.save(output_path)
    doc.close()

def process_redact_text(input_file, output_path, keyword, status_callback, stop_event):
    import fitz
    doc = fitz.open(input_file)
    total = len(doc)
    count = 0
    for i, page in enumerate(doc):
        if stop_event.is_set(): return
        status_callback(f"⬛ 正在搜尋與塗黑第 {i+1} 頁...", (i+1)/total)
        text_instances = page.search_for(keyword)
        for inst in text_instances:
            page.add_redact_annot(inst, fill=(0, 0, 0))
            count += 1
        page.apply_redactions()
        
    if count == 0:
        doc.close()
        raise Exception(f"找不到關鍵字 '{keyword}'，請確認 PDF 為可選取文字的格式。")
        
    status_callback("🔒 正在扁平化儲存確保無法復原...", 0.9)
    doc.save(output_path, deflate=True)
    doc.close()

def process_remove_watermark(input_file, output_path, status_callback, stop_event, dpi=300, position="右下角"):
    from pdf2image import convert_from_path, pdfinfo_from_path
    poppler = get_poppler_path()
    total = pdfinfo_from_path(input_file, poppler_path=poppler)["Pages"]
    temp_images = []
    
    with tempfile.TemporaryDirectory() as temp_dir:
        for i in range(1, total + 1):
            if stop_event.is_set(): break
            status_callback(f"🖌️ 正在抹除浮水印 {i} / {total}...", 0.1 + 0.7*(i/total))
            page_img = apply_watermark_removal(convert_from_path(input_file, dpi=dpi, first_page=i, last_page=i, poppler_path=poppler)[0], position)
            temp_path = os.path.join(temp_dir, f"page_{i}.jpg")
            page_img.save(temp_path, "JPEG", quality=95)
            temp_images.append(temp_path)
            del page_img; gc.collect()
            
        if stop_event.is_set(): return
        status_callback("💾 正在組合寫入檔案...", 0.9)
        
        if output_path.lower().endswith('.pptx'):
            from pptx import Presentation
            from PIL import Image
            prs = Presentation()
            with Image.open(temp_images[0]) as first_img: w_px, h_px = first_img.size
            prs.slide_width, prs.slide_height = int(w_px * 914400 / dpi), int(h_px * 914400 / dpi)
            for img_path in temp_images: prs.slides.add_slide(prs.slide_layouts[6]).shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
            prs.save(output_path)
        else:
            import fitz
            doc = fitz.open()
            for img_path in temp_images:
                img_doc = fitz.open(img_path)
                img_pdf = fitz.open("pdf", img_doc.convert_to_pdf())
                doc.insert_pdf(img_pdf)
                img_doc.close(); img_pdf.close(); del img_doc; gc.collect()
            doc.save(output_path)
            doc.close()

def process_image_ocr(input_file, output_path, status_callback, stop_event):
    from rapidocr_onnxruntime import RapidOCR
    import opencc
    img = cv2.imdecode(np.fromfile(input_file, dtype=np.uint8), cv2.IMREAD_COLOR)
    if img is None: raise Exception("無法讀取圖片。")

    status_callback("👁️ 正在進行影像強化與解析度放大...", 0.2)
    h, w = img.shape[:2]
    if min(h, w) < 1000: img = cv2.resize(img, None, fx=2.5, fy=2.5, interpolation=cv2.INTER_CUBIC)
    elif min(h, w) < 2000: img = cv2.resize(img, None, fx=1.5, fy=1.5, interpolation=cv2.INTER_CUBIC)
        
    try:
        img_yuv = cv2.cvtColor(img, cv2.COLOR_BGR2YUV)
        img_yuv[:,:,0] = cv2.createCLAHE(clipLimit=1.5, tileGridSize=(8, 8)).apply(img_yuv[:,:,0])
        img = cv2.cvtColor(img_yuv, cv2.COLOR_YUV2BGR)
    except: pass

    status_callback("👁️ 正在進行 AI 視覺文字辨識...", 0.4)
    result, _ = RapidOCR(use_angle_cls=True, det_box_thresh=0.2, text_score=0.2)(img)
    
    extracted_text = []
    if result:
        status_callback("📝 正在整理與繁化文字...", 0.8)
        converter = opencc.OpenCC('s2twp')  
        for line in result:
            if stop_event.is_set(): return
            extracted_text.append(converter.convert(line[1]))

    with open(output_path, "w", encoding="utf-8") as f: f.write("\n".join(extracted_text))
    if not result: raise Exception("找不到任何文字！請嘗試提供更清晰的圖片。")

def process_image_remove_text(input_file, output_path, status_callback, stop_event):
    from rapidocr_onnxruntime import RapidOCR
    img = cv2.imdecode(np.fromfile(input_file, dtype=np.uint8), cv2.IMREAD_COLOR)
    if img is None: raise Exception("無法讀取圖片。")

    status_callback("👁️ 正在掃描並定位圖片中的文字...", 0.3)
    result, _ = RapidOCR(use_angle_cls=True, det_box_thresh=0.2, text_score=0.2)(img)
    mask = np.zeros(img.shape[:2], dtype=np.uint8)

    if result:
        status_callback("🖌️ 正在運算遮罩與修補背景...", 0.6)
        for line in result:
            if stop_event.is_set(): return
            cv2.fillPoly(mask, [np.array(line[0], np.int32).reshape((-1, 1, 2))], 255)
        inpainted_img = cv2.inpaint(img, cv2.dilate(mask, np.ones((5, 5), np.uint8), iterations=2), inpaintRadius=5, flags=cv2.INPAINT_TELEA)
    else: inpainted_img = img 

    status_callback("💾 正在儲存無文字圖片...", 0.9)
    if not cv2.imencode(os.path.splitext(output_path)[1], inpainted_img)[1].tofile(output_path): raise Exception("儲存圖片失敗！")

def process_add_page_numbers(input_file, output_path, status_callback, stop_event):
    import fitz
    doc = fitz.open(input_file)
    total = len(doc)
    for i in range(total):
        if stop_event.is_set(): return
        status_callback(f"🔢 正在加入頁碼 ({i+1}/{total})...", (i+1)/total)
        doc[i].insert_text(fitz.Point(doc[i].rect.width / 2 - 15, doc[i].rect.height - 20), f"- {i+1} -", fontsize=11, color=(0, 0, 0))
    doc.save(output_path)
    doc.close()
