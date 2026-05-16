import os
import gc
import tempfile
import cv2
import numpy as np
from utils import get_poppler_path, apply_watermark_removal, format_size, get_model_path

# --- 基礎合併與轉換 ---

def process_merge_pdfs(input_files, output_path, status_callback, stop_event):
    from pypdf import PdfWriter
    import fitz
    import pythoncom
    import sys
    import io
    
    pythoncom.CoInitialize()
    try:
        merger = PdfWriter()
        total = len(input_files)
        with tempfile.TemporaryDirectory() as temp_dir:
            for i, file_path in enumerate(input_files):
                if stop_event.is_set(): return
                base_name = os.path.basename(file_path)
                status_callback(f"📑 處理合併: {base_name} ({i+1}/{total})", (i+1)/total)
                ext = file_path.lower().split('.')[-1]
                
                if ext == 'pdf':
                    merger.append(file_path)
                elif ext in ['jpg', 'jpeg', 'png', 'bmp']:
                    img_doc = fitz.open(file_path)
                    pdf_bytes = img_doc.convert_to_pdf()
                    img_pdf = fitz.open("pdf", pdf_bytes)
                    temp_pdf = os.path.join(temp_dir, f"temp_{i}.pdf")
                    img_pdf.save(temp_pdf)
                    img_pdf.close(); img_doc.close()
                    merger.append(temp_pdf)
                elif ext in ['docx', 'doc']:
                    from docx2pdf import convert
                    temp_pdf = os.path.join(temp_dir, f"temp_{i}.pdf")
                    dummy_out = io.StringIO()
                    old_out, old_err = sys.stdout, sys.stderr
                    sys.stdout, sys.stderr = dummy_out, dummy_out
                    try:
                        convert(file_path, temp_pdf)
                        merger.append(temp_pdf)
                    except Exception as e:
                        raise Exception(f"Word 轉換失敗 ({base_name})\n1. 請確認電腦有安裝 Microsoft Word。\n2. 請確認沒有開啟對話框卡住 Word。\n錯誤細節: {e}")
                    finally:
                        sys.stdout, sys.stderr = old_out, old_err
                else:
                    raise Exception(f"不支援的檔案格式: {ext}")
            
            if stop_event.is_set(): return
            status_callback("💾 正在輸出最終合併檔案...", 0.95)
            merger.write(output_path)
            merger.close()
    finally:
        pythoncom.CoUninitialize()

def process_images_to_pdf(input_files, output_path, status_callback, stop_event):
    import fitz  
    doc = fitz.open()
    total = len(input_files)
    for i, img_path in enumerate(input_files):
        if stop_event.is_set(): return
        status_callback(f"🖼️ 正在將圖片轉為 PDF... ({i+1}/{total})", (i+1)/total)
        img_doc = fitz.open(img_path)
        pdf_bytes = img_doc.convert_to_pdf()
        img_pdf = fitz.open("pdf", pdf_bytes)
        doc.insert_pdf(img_pdf)
        img_doc.close(); img_pdf.close()
    doc.save(output_path)
    doc.close()

# --- 加密與解鎖 ---

def process_protect_pdf(input_file, output_path, password, status_callback, stop_event):
    from pypdf import PdfWriter, PdfReader
    status_callback("🔒 正在讀取並加密 PDF...", 0.3)
    reader = PdfReader(input_file)
    writer = PdfWriter()
    for page in reader.pages:
        if stop_event.is_set(): return
        writer.add_page(page)
    status_callback("🔒 正在寫入加密檔案...", 0.7)
    writer.encrypt(password)
    with open(output_path, "wb") as f: writer.write(f)

def process_unlock_pdf(input_file, output_path, password, status_callback, stop_event):
    from pypdf import PdfWriter, PdfReader
    status_callback("🔓 正在嘗試解鎖 PDF...", 0.3)
    reader = PdfReader(input_file)
    if reader.is_encrypted:
        success = reader.decrypt(password)
        if not success: raise Exception("密碼錯誤，解鎖失敗！")
    writer = PdfWriter()
    for page in reader.pages:
        if stop_event.is_set(): return
        writer.add_page(page)
    status_callback("🔓 正在寫入無密碼檔案...", 0.7)
    with open(output_path, "wb") as f: writer.write(f)

# --- 頁面操作 ---

def parse_page_ranges(range_str, total_pages):
    if not range_str.strip(): return list(range(total_pages))
    pages = set()
    for part in range_str.replace(" ", "").split(","):
        if "-" in part:
            start, end = map(int, part.split("-"))
            pages.update(range(start - 1, end))
        else:
            pages.add(int(part) - 1)
    return sorted([p for p in pages if 0 <= p < total_pages])

def process_split_pdf(input_file, output_path, page_ranges, status_callback, stop_event):
    from pypdf import PdfWriter, PdfReader
    reader = PdfReader(input_file)
    total_pages = len(reader.pages)
    target_pages = parse_page_ranges(page_ranges, total_pages)
    
    if os.path.isdir(output_path):
        base_name = os.path.splitext(os.path.basename(input_file))[0]
        for i, p_idx in enumerate(target_pages):
            if stop_event.is_set(): return
            status_callback(f"✂️ 正在分割頁面 ({i+1}/{len(target_pages)})...", (i+1)/len(target_pages))
            writer = PdfWriter()
            writer.add_page(reader.pages[p_idx])
            with open(os.path.join(output_path, f"{base_name}_page_{p_idx+1}.pdf"), "wb") as f:
                writer.write(f)
    else:
        writer = PdfWriter()
        for i, p_idx in enumerate(target_pages):
            if stop_event.is_set(): return
            status_callback(f"✂️ 正在提取頁面 ({i+1}/{len(target_pages)})...", (i+1)/len(target_pages))
            writer.add_page(reader.pages[p_idx])
        with open(output_path, "wb") as f: writer.write(f)

def process_remove_pages(input_file, output_path, page_ranges, status_callback, stop_event):
    from pypdf import PdfWriter, PdfReader
    reader = PdfReader(input_file)
    total_pages = len(reader.pages)
    pages_to_remove = parse_page_ranges(page_ranges, total_pages)
    writer = PdfWriter()
    for i in range(total_pages):
        if stop_event.is_set(): return
        status_callback(f"🗑️ 正在掃描並剔除頁面 ({i+1}/{total_pages})...", (i+1)/total_pages)
        if i not in pages_to_remove: writer.add_page(reader.pages[i])
    with open(output_path, "wb") as f: writer.write(f)

def process_insert_blank_page(input_file, output_path, page_ranges, status_callback, stop_event):
    from pypdf import PdfWriter, PdfReader
    reader = PdfReader(input_file)
    total_pages = len(reader.pages)
    insert_after = parse_page_ranges(page_ranges, total_pages)
    writer = PdfWriter()
    for i in range(total_pages):
        if stop_event.is_set(): return
        status_callback(f"📎 正在處理並插入空白頁 ({i+1}/{total_pages})...", (i+1)/total_pages)
        page = reader.pages[i]
        writer.add_page(page)
        if i in insert_after: writer.add_blank_page(width=page.mediabox.width, height=page.mediabox.height)
    with open(output_path, "wb") as f: writer.write(f)

def process_reorder_pages(input_file, output_path, order_str, status_callback, stop_event):
    from pypdf import PdfWriter, PdfReader
    reader = PdfReader(input_file)
    writer = PdfWriter()
    indices = [int(x.strip())-1 for x in order_str.split(",") if x.strip().isdigit()]
    for i, idx in enumerate(indices):
        if stop_event.is_set(): return
        status_callback(f"🔀 正在重新排序... ({i+1}/{len(indices)})", (i+1)/len(indices))
        if 0 <= idx < len(reader.pages): writer.add_page(reader.pages[idx])
    with open(output_path, "wb") as f: writer.write(f)

def process_rotate_pdf(input_file, output_path, angle, status_callback, stop_event):
    from pypdf import PdfWriter, PdfReader
    status_callback("🔄 正在旋轉 PDF...", 0.5)
    reader = PdfReader(input_file)
    writer = PdfWriter()
    angle_int = int(angle.replace("度", ""))
    for page in reader.pages:
        if stop_event.is_set(): return
        writer.add_page(page)
        writer.pages[-1].rotate(angle_int)
    with open(output_path, "wb") as f: writer.write(f)

# --- 辦公室修改/編輯工具 ---

def process_add_page_numbers(input_file, output_path, status_callback, stop_event):
    import fitz
    doc = fitz.open(input_file)
    total = len(doc)
    for i in range(total):
        if stop_event.is_set(): return
        status_callback(f"🔢 正在加入頁碼 ({i+1}/{total})...", (i+1)/total)
        page = doc[i]
        rect = page.rect
        text = f"- {i+1} -"
        p = fitz.Point(rect.width / 2 - 15, rect.height - 20)
        page.insert_text(p, text, fontsize=11, color=(0, 0, 0))
    doc.save(output_path)
    doc.close()

def process_to_grayscale(input_file, output_path, status_callback, stop_event, dpi=200):
    import fitz
    doc = fitz.open(input_file)
    new_doc = fitz.open()
    total = len(doc)
    for i in range(total):
        if stop_event.is_set(): return
        status_callback(f"🖨️ 正在轉換為黑白/灰階 ({i+1}/{total})...", (i+1)/total)
        page = doc[i]
        pix = page.get_pixmap(dpi=dpi, colorspace=fitz.csGRAY)
        new_page = new_doc.new_page(width=page.rect.width, height=page.rect.height)
        new_page.insert_image(page.rect, pixmap=pix)
    new_doc.save(output_path, garbage=4, deflate=True)
    doc.close(); new_doc.close()

def process_flatten_pdf(input_file, output_path, status_callback, stop_event, dpi=200):
    import fitz
    doc = fitz.open(input_file)
    new_doc = fitz.open()
    total = len(doc)
    for i in range(total):
        if stop_event.is_set(): return
        status_callback(f"🥞 正在扁平化 PDF 防止篡改 ({i+1}/{total})...", (i+1)/total)
        page = doc[i]
        pix = page.get_pixmap(dpi=dpi)
        new_page = new_doc.new_page(width=page.rect.width, height=page.rect.height)
        new_page.insert_image(page.rect, pixmap=pix)
    new_doc.save(output_path, garbage=4, deflate=True)
    doc.close(); new_doc.close()

def process_extract_text(input_file, output_path, status_callback, stop_event):
    import fitz
    doc = fitz.open(input_file)
    total = len(doc)
    with open(output_path, "w", encoding="utf-8") as f:
        for i, page in enumerate(doc):
            if stop_event.is_set(): return
            status_callback(f"📄 正在提取文字 ({i+1}/{total})...", (i+1)/total)
            text = page.get_text("text")
            f.write(f"--- 第 {i+1} 頁 ---\n")
            f.write(text + "\n\n")
    doc.close()

def process_extract_original_images(input_file, output_dir, status_callback, stop_event):
    import fitz
    doc = fitz.open(input_file)
    total = len(doc)
    count = 0
    for i in range(total):
        if stop_event.is_set(): return
        status_callback(f"🖼️ 正在提取內嵌圖片 (掃描第 {i+1} 頁)...", (i+1)/total)
        for img_idx, img in enumerate(doc[i].get_images(True)):
            xref = img[0]
            base_img = doc.extract_image(xref)
            img_bytes = base_img["image"]
            ext = base_img["ext"]
            with open(os.path.join(output_dir, f"page{i+1}_img{img_idx+1}.{ext}"), "wb") as f:
                f.write(img_bytes)
            count += 1
    doc.close()
    return f"提取完成！\n共從檔案中成功提取了 {count} 張原始圖片。"

def process_compress_pdf(input_file, output_path, status_callback, stop_event):
    import fitz
    status_callback("🗜️ 正在掃描與壓縮 PDF 檔案...", 0.5)
    orig_size = os.path.getsize(input_file)
    doc = fitz.open(input_file)
    if stop_event.is_set(): return
    doc.save(output_path, garbage=4, deflate=True)
    doc.close()
    new_size = os.path.getsize(output_path)
    return f"壓縮完成！\n原大小: {format_size(orig_size)} ➡️ 新大小: {format_size(new_size)}"

# --- 浮水印處理 ---

def process_add_watermark(input_file, output_path, text, status_callback, stop_event):
    import fitz
    doc = fitz.open(input_file)
    total = len(doc)
    for i, page in enumerate(doc):
        if stop_event.is_set(): return
        status_callback(f"🖋️ 正在加入浮水印... ({i+1}/{total})", (i+1)/total)
        rect = page.rect
        page.insert_text(fitz.Point(50, rect.height / 2), text, fontsize=48, color=(1, 0, 0))
    doc.save(output_path)
    doc.close()

def process_add_image_watermark(input_file, output_path, img_path, position, status_callback, stop_event):
    import fitz
    doc = fitz.open(input_file)
    total = len(doc)
    img_doc = fitz.open(img_path)
    img_rect = img_doc[0].rect
    img_doc.close()
    for i, page in enumerate(doc):
        if stop_event.is_set(): return
        status_callback(f"🖼️ 正在壓印圖片浮水印... ({i+1}/{total})", (i+1)/total)
        page_rect = page.rect
        w = page_rect.width * 0.25
        h = w * (img_rect.height / img_rect.width)
        if position == "右下角": target_rect = fitz.Rect(page_rect.width - w - 20, page_rect.height - h - 20, page_rect.width - 20, page_rect.height - 20)
        elif position == "左下角": target_rect = fitz.Rect(20, page_rect.height - h - 20, 20 + w, page_rect.height - 20)
        elif position == "右上角": target_rect = fitz.Rect(page_rect.width - w - 20, 20, page_rect.width - 20, 20 + h)
        else: target_rect = fitz.Rect(20, 20, 20 + w, 20 + h)
        page.insert_image(target_rect, filename=img_path)
    doc.save(output_path)
    doc.close()

def process_remove_watermark(input_file, output_path, status_callback, stop_event, dpi=300, position="右下角"):
    from pdf2image import convert_from_path, pdfinfo_from_path
    poppler = get_poppler_path()
    info = pdfinfo_from_path(input_file, poppler_path=poppler)
    total = info["Pages"]
    temp_images = []
    
    with tempfile.TemporaryDirectory() as temp_dir:
        for i in range(1, total + 1):
            if stop_event.is_set(): break
            status_callback(f"🖌️ 正在抹除浮水印 {i} / {total}...", 0.1 + 0.7*(i/total))
            page_img = convert_from_path(input_file, dpi=dpi, first_page=i, last_page=i, poppler_path=poppler)[0]
            page_img = apply_watermark_removal(page_img, position)
            
            temp_path = os.path.join(temp_dir, f"page_{i}.jpg")
            page_img.save(temp_path, "JPEG", quality=95)
            temp_images.append(temp_path)
            del page_img; gc.collect()
            
        if stop_event.is_set(): return
        status_callback("💾 正在組合寫入檔案...", 0.9)
        
        if output_path.lower().endswith('.pptx'):
            from pptx import Presentation
            from PIL import Image
            prs = Presentation()
            with Image.open(temp_images[0]) as first_img: width_px, height_px = first_img.size
            prs.slide_width = int(width_px * 914400 / dpi) 
            prs.slide_height = int(height_px * 914400 / dpi)
            for img_path in temp_images:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
            prs.save(output_path)
        else:
            import fitz
            doc = fitz.open()
            for img_path in temp_images:
                img_doc = fitz.open(img_path)
                pdf_bytes = img_doc.convert_to_pdf()
                img_pdf = fitz.open("pdf", pdf_bytes)
                doc.insert_pdf(img_pdf)
                img_doc.close(); img_pdf.close()
            doc.save(output_path)
            doc.close()

def process_pdf_to_images(input_file, output_dir, status_callback, stop_event, dpi=300):
    from pdf2image import convert_from_path, pdfinfo_from_path
    poppler = get_poppler_path()
    info = pdfinfo_from_path(input_file, poppler_path=poppler)
    total = info["Pages"]
    base_name = os.path.splitext(os.path.basename(input_file))[0]
    for i in range(1, total + 1):
        if stop_event.is_set(): break
        status_callback(f"🖼️ 正在處理並儲存圖片 {i} / {total}...", i/total)
        page_img = convert_from_path(input_file, dpi=dpi, first_page=i, last_page=i, poppler_path=poppler)[0]
        page_img.save(os.path.join(output_dir, f"{base_name}_{i}.jpg"), 'JPEG')
        del page_img; gc.collect()

# --- 高精度 PDF/圖片 轉 PPT (雙軌辨識 + 表格原生還原 + 格式還原) ---

def process_pdf_to_ppt(input_file, output_path, status_callback, stop_event, dpi=300, ppt_mode="圖文排版 (智慧 OCR)"):
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    import fitz
    from PIL import Image, ImageDraw
    import tempfile
    import os
    from collections import Counter
    import opencc
    from bs4 import BeautifulSoup
    import cv2
    import numpy as np
    
    # 強制繁化
    converter = opencc.OpenCC('s2twp')

    is_pdf = input_file.lower().endswith('.pdf')
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6] 
    
    def expand_bbox(x0, y0, x1, y1, scale=1.1, max_w=9999, max_h=9999):
        cx, cy = (x0 + x1) / 2, (y0 + y1) / 2
        w, h = (x1 - x0) * scale, (y1 - y0) * scale
        new_x0 = max(0, cx - w / 2)
        new_y0 = max(0, cy - h / 2)
        new_x1 = min(max_w, cx + w / 2)
        new_y1 = min(max_h, cy + h / 2)
        return [new_x0, new_y0, new_x1, new_y1]

    def get_dynamic_bg_color(img_obj, px_bbox):
        w, h = img_obj.size
        x0, y0, x1, y1 = [int(v) for v in px_bbox]
        samples = []
        offsets = [
            (x0 - 5, y0 - 5), (x0 + (x1-x0)//2, y0 - 5), (x1 + 5, y0 - 5),
            (x0 - 5, y1 + 5), (x0 + (x1-x0)//2, y1 + 5), (x1 + 5, y1 + 5),
            (x0 - 5, y0 + (y1-y0)//2), (x1 + 5, y0 + (y1-y0)//2)
        ]
        for sx, sy in offsets:
            sx = max(0, min(w - 1, sx))
            sy = max(0, min(h - 1, sy))
            try:
                pixel = img_obj.getpixel((sx, sy))
                if isinstance(pixel, int): pixel = (pixel, pixel, pixel)
                samples.append(pixel[:3]) 
            except: pass
        return Counter(samples).most_common(1)[0][0] if samples else (255, 255, 255)

    def draw_ppt_table_from_html(slide, html_str, x_pt, y_pt, w_pt, h_pt):
        soup = BeautifulSoup(html_str, 'html.parser')
        rows = soup.find_all('tr')
        if not rows: return
        
        num_rows = len(rows)
        num_cols = max(len(row.find_all(['td', 'th'])) for row in rows)
        
        table_shape = slide.shapes.add_table(num_rows, num_cols, Pt(x_pt), Pt(y_pt), Pt(w_pt), Pt(h_pt))
        table = table_shape.table
        
        for r_idx, row in enumerate(rows):
            cols = row.find_all(['td', 'th'])
            for c_idx, col in enumerate(cols):
                if r_idx < len(table.rows) and c_idx < len(table.columns):
                    cell = table.cell(r_idx, c_idx)
                    text = converter.convert(col.get_text(strip=True))
                    cell.text = text
                    
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = "微軟正黑體"
                            run.font.size = Pt(10)
                            run.font.color.rgb = RGBColor(0, 0, 0)

    # 動態路徑載入模型
    ocr_engine, layout_engine, table_engine = None, None, None
    if ppt_mode == "圖文排版 (智慧 OCR)":
        from rapidocr_onnxruntime import RapidOCR
        from rapid_layout import RapidLayout
        from rapid_table import RapidTable
        from utils import get_model_path
        
        model_dir = get_model_path()
        det_path = os.path.join(model_dir, "ch_PP-OCRv4_det_infer.onnx")
        cls_path = os.path.join(model_dir, "ch_ppocr_mobile_v2.0_cls_infer.onnx")
        rec_path = os.path.join(model_dir, "ch_PP-OCRv4_rec_infer.onnx")
        layout_path = os.path.join(model_dir, "layout_cdla_dataset.onnx")
        table_path = os.path.join(model_dir, "ch_ppstructure_mobile_v2.0_SLANet.onnx")
        
        if os.path.exists(det_path):
            ocr_engine = RapidOCR(det_model_path=det_path, cls_model_path=cls_path, rec_model_path=rec_path, det_box_thresh=0.3)
            layout_engine = RapidLayout(model_path=layout_path)
            table_engine = RapidTable(model_path=table_path)
        else:
            ocr_engine = RapidOCR()
            layout_engine = RapidLayout()
            table_engine = RapidTable()

    with tempfile.TemporaryDirectory() as temp_dir:
        def process_slide(hr_img_path, normal_img_path, slide_obj, scale_down):
            img_obj = Image.open(normal_img_path).convert("RGB")
            draw = ImageDraw.Draw(img_obj)
            text_boxes_data = []
            table_html_data = []
            
            # 解決 Windows 中文路徑報錯
            hr_img = cv2.imdecode(np.fromfile(hr_img_path, dtype=np.uint8), cv2.IMREAD_COLOR)
            if hr_img is None:
                return 

            # AI 視覺預處理模組 (提升辨識率)
            def enhance_for_ocr(image):
                try:
                    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
                    clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
                    enhanced = clahe.apply(gray)
                    kernel = np.array([[0, -0.5, 0], [-0.5, 3, -0.5], [0, -0.5, 0]])
                    sharpened = cv2.filter2D(enhanced, -1, kernel)
                    return cv2.cvtColor(sharpened, cv2.COLOR_GRAY2BGR)
                except:
                    return image 

            raw_layout = layout_engine(hr_img)
            regions = []
            if isinstance(raw_layout, tuple): raw_layout = raw_layout[0]
            
            if isinstance(raw_layout, list):
                for item in raw_layout:
                    if isinstance(item, dict): regions.append((item.get('bbox'), item.get('label')))
            elif hasattr(raw_layout, 'boxes') and hasattr(raw_layout, 'class_names'):
                for box, label in zip(raw_layout.boxes, raw_layout.class_names):
                    regions.append((box, label))

            if regions:
                for box, label in regions:
                    if box is None or label is None: continue
                    if len(box) == 4 and not isinstance(box[0], (list, tuple, np.ndarray)):
                        bx0, by0, bx1, by1 = box
                    else:
                        xs = [p[0] for p in box]; ys = [p[1] for p in box]
                        bx0, by0, bx1, by1 = min(xs), min(ys), max(xs), max(ys)
                        
                    x0, y0, x1, y1 = [v * scale_down for v in [bx0, by0, bx1, by1]]
                    
                    if label == 'table':
                        table_crop = hr_img[int(by0):int(by1), int(bx0):int(bx1)]
                        if table_crop is not None and table_crop.size > 0:
                            enhanced_table = enhance_for_ocr(table_crop)
                            raw_table = table_engine(enhanced_table)
                            table_res = raw_table[0] if isinstance(raw_table, tuple) else raw_table
                            
                            html_str = None
                            if isinstance(table_res, str): html_str = table_res
                            elif isinstance(table_res, dict) and 'html' in table_res: html_str = table_res['html']
                            elif hasattr(table_res, 'html'): html_str = table_res.html
                            
                            if html_str:
                                table_html_data.append({"html": html_str, "x": x0, "y": y0, "w": x1-x0, "h": y1-y0})
                                exp_px_bbox = expand_bbox(x0, y0, x1, y1, scale=1.02, max_w=img_obj.width, max_h=img_obj.height)
                                bg_color = get_dynamic_bg_color(img_obj, exp_px_bbox)
                                draw.rectangle(exp_px_bbox, fill=bg_color)
                                
                    elif label in ['text', 'title', 'figure']:
                        text_crop = hr_img[int(by0):int(by1), int(bx0):int(bx1)]
                        if text_crop is not None and text_crop.size > 0:
                            enhanced_text = enhance_for_ocr(text_crop)
                            raw_ocr = ocr_engine(enhanced_text)
                            ocr_res = raw_ocr[0] if isinstance(raw_ocr, tuple) else raw_ocr
                            if ocr_res:
                                for line in ocr_res:
                                    line_box = line[0]
                                    text = converter.convert(line[1])
                                    lx0 = min([p[0] for p in line_box]) * scale_down + x0
                                    ly0 = min([p[1] for p in line_box]) * scale_down + y0
                                    lx1 = max([p[0] for p in line_box]) * scale_down + x0
                                    ly1 = max([p[1] for p in line_box]) * scale_down + y0
                                    
                                    text_boxes_data.append({"text": text, "x": lx0, "y": ly0, "w": lx1-lx0, "h": ly1-ly0})
                                    exp_px_bbox = expand_bbox(lx0, ly0, lx1, ly1, scale=1.1, max_w=img_obj.width, max_h=img_obj.height)
                                    bg_color = get_dynamic_bg_color(img_obj, exp_px_bbox)
                                    draw.rectangle(exp_px_bbox, fill=bg_color)
            else:
                enhanced_full_img = enhance_for_ocr(hr_img)
                raw_ocr = ocr_engine(enhanced_full_img)
                ocr_res = raw_ocr[0] if isinstance(raw_ocr, tuple) else raw_ocr
                if ocr_res:
                    for line in ocr_res:
                        line_box = line[0]
                        text = converter.convert(line[1])
                        lx0 = min([p[0] for p in line_box]) * scale_down
                        ly0 = min([p[1] for p in line_box]) * scale_down
                        lx1 = max([p[0] for p in line_box]) * scale_down
                        ly1 = max([p[1] for p in line_box]) * scale_down
                        
                        text_boxes_data.append({"text": text, "x": lx0, "y": ly0, "w": lx1-lx0, "h": ly1-ly0})
                        exp_px_bbox = expand_bbox(lx0, ly0, lx1, ly1, scale=1.1, max_w=img_obj.width, max_h=img_obj.height)
                        bg_color = get_dynamic_bg_color(img_obj, exp_px_bbox)
                        draw.rectangle(exp_px_bbox, fill=bg_color)

            img_obj.save(normal_img_path, "JPEG", quality=90)
            slide_obj.shapes.add_picture(normal_img_path, 0, 0, prs.slide_width, prs.slide_height)
            
            for item in text_boxes_data:
                txBox = slide_obj.shapes.add_textbox(Pt(item["x"]), Pt(item["y"]), Pt(item["w"]), Pt(item["h"]))
                tf = txBox.text_frame
                tf.clear()
                tf.word_wrap = False
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = item["text"]
                run.font.size = Pt(max(8, item["h"] * 0.75))
                run.font.name = "微軟正黑體"
                run.font.color.rgb = RGBColor(0, 0, 0)
            
            for t_item in table_html_data:
                draw_ppt_table_from_html(slide_obj, t_item["html"], t_item["x"], t_item["y"], t_item["w"], t_item["h"])

        if is_pdf:
            doc = fitz.open(input_file)
            total = len(doc)
            for i in range(total):
                if stop_event.is_set(): break
                page = doc[i]
                if i == 0:
                    prs.slide_width = Pt(page.rect.width)
                    prs.slide_height = Pt(page.rect.height)
                
                slide = prs.slides.add_slide(blank_slide_layout)
                normal_path = os.path.join(temp_dir, f"bg_{i}.jpg")
                
                if ppt_mode == "純圖片簡報 (較快)":
                    status_callback(f"📊 正在轉換純圖 PPT (第 {i+1} / {total} 頁)...", (i+1)/total)
                    page.get_pixmap(dpi=dpi, colorspace=fitz.csRGB).save(normal_path)
                    slide.shapes.add_picture(normal_path, 0, 0, prs.slide_width, prs.slide_height)
                    continue
                
                status_callback(f"👁️ 進行版面分析與雙軌重建 (第 {i+1} / {total} 頁)...", (i+1)/total)
                
                mat = fitz.Matrix(4.0, 4.0)
                hr_path = os.path.join(temp_dir, f"hr_{i}.jpg")
                page.get_pixmap(matrix=mat, colorspace=fitz.csRGB).save(hr_path)
                page.get_pixmap(dpi=72, colorspace=fitz.csRGB).save(normal_path)
                
                process_slide(hr_path, normal_path, slide, scale_down=1.0/4.0)
            doc.close()
        else:
            img = Image.open(input_file).convert('RGB')
            normal_path = os.path.join(temp_dir, "normal.jpg")
            hr_path = os.path.join(temp_dir, "hr.jpg")
            
            img.save(normal_path, "JPEG", quality=95)
            hr_img = img.resize((img.width * 4, img.height * 4), Image.LANCZOS)
            hr_img.save(hr_path, "JPEG", quality=100)
            
            prs.slide_width = Pt(img.width * 72 / dpi)
            prs.slide_height = Pt(img.height * 72 / dpi)
            slide = prs.slides.add_slide(blank_slide_layout)

            if ppt_mode == "純圖片簡報 (較快)":
                status_callback("🖼️ 正在處理圖片檔案...", 0.5)
                slide.shapes.add_picture(normal_path, 0, 0, prs.slide_width, prs.slide_height)
            else:
                status_callback("👁️ 進行版面分析與雙軌重建...", 0.5)
                process_slide(hr_path, normal_path, slide, scale_down=1.0/4.0)

        if not stop_event.is_set():
            status_callback("💾 正在儲存檔案...", 0.95)
            prs.save(output_path)

# --- 圖片進階處理 (OCR & 文字抹除) ---

def process_image_ocr(input_file, output_path, status_callback, stop_event):
    from rapidocr_onnxruntime import RapidOCR
    import opencc
    import cv2
    import numpy as np
    
    converter = opencc.OpenCC('s2twp')  # 轉換為繁體中文
    
    # 支援讀取包含中文路徑的圖片
    img = cv2.imdecode(np.fromfile(input_file, dtype=np.uint8), cv2.IMREAD_COLOR)
    if img is None:
        raise Exception("無法讀取圖片，請確認檔案是否損壞。")

    status_callback("👁️ 正在進行 AI 視覺文字辨識...", 0.4)
    ocr_engine = RapidOCR()
    result, _ = ocr_engine(img)
    
    extracted_text = []
    if result:
        status_callback("📝 正在整理與繁化文字...", 0.8)
        for line in result:
            if stop_event.is_set(): return
            text = converter.convert(line[1])
            extracted_text.append(text)

    with open(output_path, "w", encoding="utf-8") as f:
        f.write("\n".join(extracted_text))
        
    if not result:
        raise Exception("找不到任何文字！")

def process_image_remove_text(input_file, output_path, status_callback, stop_event):
    from rapidocr_onnxruntime import RapidOCR
    import cv2
    import numpy as np

    status_callback("👁️ 正在掃描並定位圖片中的文字...", 0.3)
    img = cv2.imdecode(np.fromfile(input_file, dtype=np.uint8), cv2.IMREAD_COLOR)
    if img is None:
        raise Exception("無法讀取圖片，請確認檔案是否損壞。")

    ocr_engine = RapidOCR()
    result, _ = ocr_engine(img)
    
    # 建立純黑遮罩
    mask = np.zeros(img.shape[:2], dtype=np.uint8)

    if result:
        status_callback("🖌️ 正在運算遮罩與修補背景...", 0.6)
        for line in result:
            if stop_event.is_set(): return
            box = line[0]
            # 將浮點數座標轉為整數座標矩陣
            pts = np.array(box, np.int32).reshape((-1, 1, 2))
            # 在遮罩上畫出白色多邊形 (標示文字範圍)
            cv2.fillPoly(mask, [pts], 255)

        # 稍微膨脹遮罩，確保文字邊緣被完全覆蓋
        kernel = np.ones((5, 5), np.uint8)
        mask = cv2.dilate(mask, kernel, iterations=2)

        # 使用 Telea 演算法進行影像修補 (Inpainting)
        inpainted_img = cv2.inpaint(img, mask, inpaintRadius=5, flags=cv2.INPAINT_TELEA)
    else:
        # 如果沒發現文字，直接輸出原圖
        inpainted_img = img 

    status_callback("💾 正在儲存無文字圖片...", 0.9)
    # 支援存入包含中文路徑的資料夾
    ext = os.path.splitext(output_path)[1]
    is_success, im_buf_arr = cv2.imencode(ext, inpainted_img)
    if is_success:
        im_buf_arr.tofile(output_path)
    else:
        raise Exception("儲存圖片失敗！")
