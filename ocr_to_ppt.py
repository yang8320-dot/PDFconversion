import os
import gc
import easyocr
import fitz
from pdf2image import convert_from_path, pdfinfo_from_path
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from utils import get_poppler_path, get_model_path
from PIL import Image

def process_ocr_to_ppt(input_file, output_path, status_callback, stop_event, use_gpu, white_bg, use_ocr, dpi):
    is_pdf = input_file.lower().endswith('.pdf')
    prs = Presentation()
    
    # ==========================================
    # === 智慧原生還原模式 (極近商業軟體效果) ===
    # ==========================================
    if not use_ocr and is_pdf:
        status_callback("⚡ 啟動高階原生還原引擎 (無疊影模式)...", 0.1)
        doc = fitz.open(input_file)
        total_pages = len(doc)
        
        for i in range(total_pages):
            if stop_event.is_set(): break
            status_callback(f"⚡ 正在精準重建 PPT (第 {i+1} / {total_pages} 頁)...", 0.1 + 0.8 * ((i+1)/total_pages))
            
            page = doc[i]
            
            # 1. 抓取所有原生文字區塊 (Dict 格式包含字體、大小、顏色與座標)
            text_data = page.get_text("dict")["blocks"]
            
            # 2. 【魔法橡皮擦】：在把頁面轉成背景圖片前，將所有文字用「白色方塊」塗掉！
            # 這樣背景圖就只會剩下線條、圖片、表格框線，不會有原來的字
            for b in text_data:
                if b["type"] == 0:  # 0代表文字區塊
                    page.draw_rect(b["bbox"], color=(1, 1, 1), fill=(1, 1, 1)) # 畫白底覆蓋
            
            # 3. 輸出乾淨無字的背景圖
            pix = page.get_pixmap(dpi=dpi)
            temp_img = f"temp_fast_bg_{i}.jpg"
            pix.save(temp_img)
            
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(temp_img, 0, 0, width=prs.slide_width, height=prs.slide_height)
            os.remove(temp_img)

            # 4. 【Span-Level 精準排版】：還原字體大小、顏色與絕對座標
            scale_w = prs.slide_width / page.rect.width
            scale_h = prs.slide_height / page.rect.height
            
            for b in text_data:
                if b["type"] == 0:
                    for line in b["lines"]:
                        for span in line["spans"]: # Span 是最小的文字單位，確保精準
                            text = span["text"].strip()
                            if not text: continue
                            
                            x0, y0, x1, y1 = span["bbox"]
                            left, top = int(x0 * scale_w), int(y0 * scale_h)
                            width, height = int((x1 - x0) * scale_w), int((y1 - y0) * scale_h)
                            
                            # 若寬高異常，給予防呆機制
                            if width <= 0 or height <= 0: continue
                            
                            txBox = slide.shapes.add_textbox(left, top, width, height)
                            tf = txBox.text_frame
                            tf.word_wrap = False
                            tf.vertical_anchor = MSO_ANCHOR.TOP
                            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
                            
                            p = tf.add_paragraph()
                            p.text = text
                            
                            # 讀取並設定 PDF 原生字體大小
                            origin_size = span["size"]
                            p.font.size = Pt(max(6, min(int(origin_size), 96)))
                            
                            # 讀取並設定 PDF 原生文字顏色
                            color_int = span["color"]
                            r = (color_int >> 16) & 255
                            g = (color_int >> 8) & 255
                            b_col = color_int & 255
                            p.font.color.rgb = RGBColor(r, g, b_col)
                            
                            # 統一使用微軟正黑體確保中文顯示正常
                            p.font.name = '微軟正黑體'
                            
                            # 若字體為粗體，進行還原
                            if "bold" in span["font"].lower():
                                p.font.bold = True
        doc.close()
    
    # ==========================================
    # === 傳統 OCR 圖片辨識模式 (用於掃描檔) ===
    # ==========================================
    else:
        status_callback("⏳ 正在初始化 OCR 模型 (圖片強行辨識)...", 0.05)
        reader = easyocr.Reader(['ch_tra', 'en'], model_storage_directory=get_model_path(), download_enabled=False, gpu=use_gpu)
        poppler = get_poppler_path()
        
        if is_pdf:
            total_pages = pdfinfo_from_path(input_file, poppler_path=poppler)["Pages"]
        else:
            img_cache = [Image.open(input_file).convert('RGB')]
            total_pages = 1

        for i in range(total_pages):
            if stop_event.is_set(): break
            status_callback(f"🔍 正在辨識與排版 PPT (第 {i+1} / {total_pages} 頁)...", 0.1 + 0.8 * ((i+1)/total_pages))
            
            if is_pdf: page_img = convert_from_path(input_file, dpi=dpi, first_page=i+1, last_page=i+1, poppler_path=poppler)[0]
            else: page_img = img_cache[0]
                
            img_width, img_height = page_img.size
            temp_img = f"temp_ppt_page_{i}.jpg"
            page_img.save(temp_img, 'JPEG')
            
            result = reader.readtext(temp_img)
            slide = prs.slides.add_slide(prs.slide_layouts[6]) 
            slide.shapes.add_picture(temp_img, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
            for (bbox, text, prob) in result:
                x_tl, y_tl = bbox[0]; x_br, y_br = bbox[2]
                left = int((x_tl / img_width) * prs.slide_width)
                top = int((y_tl / img_height) * prs.slide_height)
                width = int(((x_br - x_tl) / img_width) * prs.slide_width)
                height = int(((y_br - y_tl) / img_height) * prs.slide_height)
                
                txBox = slide.shapes.add_textbox(left, top, width, height)
                
                # 若為 OCR 模式且勾選白底，則上白底蓋住原字
                if white_bg:
                    fill = txBox.fill; fill.solid(); fill.fore_color.rgb = RGBColor(255, 255, 255)
                    
                tf = txBox.text_frame
                tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.TOP
                tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
                p = tf.add_paragraph(); p.text = text
                estimated_pt = max(6, min(int((height / 12700) * 0.75), 96))
                p.font.size = Pt(estimated_pt); p.font.name = '微軟正黑體'
                    
            os.remove(temp_img)
            if is_pdf: del page_img; gc.collect()
    
    if not stop_event.is_set():
        status_callback("💾 正在寫入 PPT 檔案...", 0.95)
        prs.save(output_path)
