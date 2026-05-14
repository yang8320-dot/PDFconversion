import os
import sys
import ctypes
import threading
import tkinter as tk
from tkinter import filedialog, messagebox
from tkinterdnd2 import DND_FILES, TkinterDnD
import easyocr
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Pt

# ==========================================
# 1. 程式防重啟 (Single Instance Check)
# ==========================================
def check_single_instance():
    mutex_name = "Global\\PDF_TO_PPT_OCR_MUTEX"
    kernel32 = ctypes.windll.kernel32
    # 建立系統層級的互斥鎖
    mutex = kernel32.CreateMutexW(None, False, mutex_name)
    last_error = kernel32.GetLastError()
    
    if last_error == 183: # ERROR_ALREADY_EXISTS (183)
        messagebox.showwarning("提示", "程式已經在執行中，請勿重複開啟。")
        sys.exit(0)
    return mutex # 必須回傳並保持參考，否則會被垃圾回收機制清除

# ==========================================
# 2. 介面 DPI 修正 (DPI Awareness)
# ==========================================
def set_dpi_awareness():
    try:
        # 告訴 Windows 10/11 程式支援高解析度縮放
        ctypes.windll.shcore.SetProcessDpiAwareness(2) # PROCESS_PER_MONITOR_DPI_AWARE
    except Exception:
        try:
            # 舊版 Windows 的備用方案
            ctypes.windll.user32.SetProcessDPIAware()
        except Exception:
            pass

# ==========================================
# 核心轉換邏輯
# ==========================================
def get_poppler_path():
    if getattr(sys, 'frozen', False):
        base_dir = os.path.dirname(sys.executable)
        return os.path.join(base_dir, "Library", "poppler_bin")
    else:
        return os.path.join(os.path.abspath("."), "Library", "poppler_bin")

class PDF2PPTApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PDF轉PPT OCR 工具")
        self.root.geometry("600x400")
        self.root.configure(bg="#f0f0f0")
        
        # 建立拖曳與點擊區域
        self.drop_frame = tk.Frame(self.root, bg="#ffffff", bd=2, relief="groove")
        self.drop_frame.pack(expand=True, fill=tk.BOTH, padx=30, pady=30)
        
        self.status_label = tk.Label(
            self.drop_frame, 
            text="📁 將 PDF 檔案拖曳至此\n或\n點擊這裡選擇檔案", 
            bg="#ffffff", fg="#555555", font=("Microsoft JhengHei", 14, "bold"),
            justify="center"
        )
        self.status_label.pack(expand=True)
        
        # 綁定滑鼠點擊事件
        self.drop_frame.bind("<Button-1>", lambda e: self.browse_file())
        self.status_label.bind("<Button-1>", lambda e: self.browse_file())
        
        # 綁定拖曳事件 (Drag & Drop)
        self.root.drop_target_register(DND_FILES)
        self.root.dnd_bind('<<Drop>>', self.on_drop)

    def on_drop(self, event):
        # 處理 tkinterdnd2 傳回路徑可能包含 {} 的問題
        files = self.root.tk.splitlist(event.data)
        if not files: return
        file_path = files[0]
        
        if not file_path.lower().endswith('.pdf'):
            messagebox.showerror("錯誤", "請提供有效的 PDF 檔案！")
            return
            
        self.ask_save_and_process(file_path)

    def browse_file(self):
        file_path = filedialog.askopenfilename(
            title="選擇要轉換的 PDF 檔案",
            filetypes=[("PDF Files", "*.pdf")]
        )
        if file_path:
            self.ask_save_and_process(file_path)

    def ask_save_and_process(self, pdf_path):
        ppt_path = filedialog.asksaveasfilename(
            title="選擇 PPT 儲存位置",
            defaultextension=".pptx",
            filetypes=[("PowerPoint Files", "*.pptx")]
        )
        if not ppt_path:
            return
            
        # 啟動背景執行緒處理 OCR，避免 GUI 視窗凍結 (Not Responding)
        threading.Thread(target=self.run_ocr_task, args=(pdf_path, ppt_path), daemon=True).start()

    def update_status(self, text):
        # 安全地在主執行緒更新 UI
        self.root.after(0, lambda: self.status_label.config(text=text))

    def run_ocr_task(self, pdf_path, ppt_output_path):
        try:
            self.update_status("⏳ 正在初始化 OCR 模型...")
            # 關閉視窗的互動能力，避免轉換途中重複拖曳
            self.drop_frame.unbind("<Button-1>")
            self.status_label.unbind("<Button-1>")
            self.root.dnd_bind('<<Drop>>', '')

            languages = ['ch_tra', 'ch_sim', 'en']
            reader = easyocr.Reader(languages)
            
            poppler_path = get_poppler_path()
            self.update_status("📄 正在將 PDF 轉換為圖片...")
            pages = convert_from_path(pdf_path, dpi=200, poppler_path=poppler_path)
            
            prs = Presentation()
            total_pages = len(pages)
            
            for i, page_img in enumerate(pages):
                self.update_status(f"🔍 正在辨識第 {i+1} / {total_pages} 頁...")
                temp_img = f"temp_page_{i}.jpg"
                page_img.save(temp_img, 'JPEG')
                
                result = reader.readtext(temp_img, detail=0)
                extracted_text = "\n".join(result)
                
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), prs.slide_height - Inches(1))
                txBox.text_frame.word_wrap = True
                p = txBox.text_frame.add_paragraph()
                p.text = extracted_text
                p.font.size = Pt(14)
                
                os.remove(temp_img)
            
            self.update_status("💾 正在儲存 PPT...")
            prs.save(ppt_output_path)
            
            self.update_status("✅ 轉換成功！\n可以繼續拖曳下一個檔案")
            messagebox.showinfo("完成", f"轉換成功！\n檔案已儲存至：\n{ppt_output_path}")

        except Exception as e:
            self.update_status("❌ 發生錯誤，請重試")
            messagebox.showerror("錯誤", f"轉換過程中發生錯誤：\n{e}")
        finally:
            # 恢復互動能力
            self.drop_frame.bind("<Button-1>", lambda e: self.browse_file())
            self.status_label.bind("<Button-1>", lambda e: self.browse_file())
            self.root.dnd_bind('<<Drop>>', self.on_drop)

def main():
    # 1. DPI 修正
    set_dpi_awareness()
    
    # 2. 防重啟檢查
    _mutex = check_single_instance()
    
    # 3. 啟動 TkinterDnD (支援拖曳的 Tk 視窗)
    root = TkinterDnD.Tk()
    app = PDF2PPTApp(root)
    root.mainloop()

if __name__ == "__main__":
    main()
